from flask import Flask, render_template, request, jsonify, send_file
import os
import re
import random
from pptx import Presentation
from pptx.util import Inches
from dotenv import load_dotenv
from openai import AzureOpenAI  # Import AzureOpenAI from the openai package

# Load environment variables from .env file
load_dotenv()

# Configuration for Azure OpenAI
AZURE_OPENAI_API_KEY = os.getenv("AZURE_OPENAI_API_KEY")
AZURE_OPENAI_ENDPOINT = os.getenv("AZURE_OPENAI_ENDPOINT")
AZURE_OPENAI_DEPLOYMENT_ID = os.getenv('AZURE_OPENAI_DEPLOYMENT_ID')
MODEL_NAME = "gpt-4o"  # Your model name

app = Flask(__name__)

# Create Azure OpenAI client
client = AzureOpenAI(
    azure_endpoint=AZURE_OPENAI_ENDPOINT,
    api_version="2024-02-01",
    api_key=AZURE_OPENAI_API_KEY
)

def generate_presentation_content(n, topic, description):
    # Construct the prompt for Azure OpenAI API
    prompt = f"""
    You are an AI assistant that helps users create presentations. 
    Generate a well detailed and informative {n}-slide presentation in german language on the topic "{topic}" with the description: "{description}". 
    Each slide should have the following format:
    
    **Slide X:**
    **Title:** [Insert Title]
    **Key Points:**
    - Bullet point 1
    - Bullet point 2
    - Bullet point 3
    """

    # Making a request to the Azure OpenAI API
    completion = client.chat.completions.create(
        model=MODEL_NAME,
        messages=[{"role": "user", "content": prompt}],
    )

    # Check for the response
    return completion.choices[0].message.content

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/generate', methods=['POST'])
def generate():
    try:
        # Get data from the form
        data = request.get_json()
        n = data['slides']
        topic = data['topic']
        description = data['description']

        # Generate presentation content using Azure OpenAI
        generated_content = generate_presentation_content(n, topic, description)

        # Return the generated content as JSON
        return jsonify({'content': generated_content})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/generate_ppt', methods=['POST'])
def generate_ppt():
    # Get the JSON data from the request
    data = request.get_json()

    if 'content' not in data:
        return jsonify({"error": "No content provided."}), 400

    content = data['content']

    # Select a random template from the 'temps' folder
    template_folder = 'temps'
    template_files = [f for f in os.listdir(template_folder) if f.endswith('.pptx')]

    if not template_files:
        return jsonify({"error": "No template files found."}), 500

    # Randomly select a template
    selected_template = random.choice(template_files)
    template_ppt_path = os.path.join(template_folder, selected_template)

    # Load the template presentation
    presentation = Presentation(template_ppt_path)

    # Split the content into slides based on double newlines
    slides_content = re.split(r'\n\s*\n+', content.strip())

    for i, slide_content in enumerate(slides_content):
        title = ""
        subtitle = ""
        main_text = ""

        # Extract title and key points
        title_match = re.search(r'\*\*Title:\*\* (.+)', slide_content)
        key_points_match = re.findall(r'- (.+)', slide_content)

        if title_match:
            title = title_match.group(1).strip()

        # Skip creating a slide if the title is empty
        if not title:
            continue

        # Reuse existing slides or add new ones if the template doesn't have enough
        if i < len(presentation.slides):
            slide = presentation.slides[i]
        else:
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Use a blank layout

        # Set the title
        if slide.shapes.title:
            slide.shapes.title.text = title

        # Add key points with bullet points
        if key_points_match:
            left = Inches(1)
            top = Inches(3)
            width = Inches(8)
            height = Inches(3)
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            for key_point in key_points_match:
                p = text_frame.add_paragraph()
                p.text = f"• {key_point.strip()}"
                p.space_after = Inches(0.1)
                p.font.size = Inches(0.25)
                p.level = 0  # This will keep the bullet points aligned at the first level

    # Save the presentation
    ppt_file_path = 'presentation.pptx'
    presentation.save(ppt_file_path)

    # Return the file for download
    return send_file(ppt_file_path, as_attachment=True, download_name='presentation.pptx')

if __name__ == '__main__':
    app.run(debug=True)
